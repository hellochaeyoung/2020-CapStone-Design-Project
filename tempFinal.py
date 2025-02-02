import cv2
import numpy as np
import scipy as sp
import math
from math import *
from PIL import Image, ImageFilter
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Mm
from pptx.util import Cm
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
#import matplotlib.pyplot as plt
import matplotlib as mpl
import pytesseract


class Yolo:
    
    width = 0
    height = 0
    net = None
    classes = []
    output_layers = []
    colors = np.random.uniform(0,255,3) # 단순 선언
    roi = 0
    image = None
    presentation = None
    
    
    # Load Yolo
    def __init__(self):
        print("yolo init ok")
#        self.net = cv2.dnn.readNet("C:/Users/hansung/darkflow/bin/yolov4_best.weights", "C:/Users/hansung/darkflow/cfg/yolov4.cfg") #// 네트워크 로드
        self.net = cv2.dnn.readNet("C:/Users/hansung/darkflow/bin/yolov3_best.weights", "C:/Users/hansung/darkflow/cfg/yolov3.cfg") #// 네트워크 로드

        #classes = []

        with open("C:/Users/hansung/darkflow/cfg/obj.names", "r") as f: 
            self.classes = [line.strip() for line in f.readlines()] 
        layer_names = self.net.getLayerNames()
        self.output_layers = [layer_names[i[0] - 1] for i in self.net.getUnconnectedOutLayers()]
        self.colors = np.random.uniform(0, 255, size=(len(self.classes), 3))

    # Loading image
    def loadImage(self, image):
        self.height, self.width, channels = image.shape
        self.image = image
        
        return image
        

    # Detecting objects
    def detectingObjects(self,image):
        print("detectingObjects ok")
        blob = cv2.dnn.blobFromImage(image, 0.00392, (416, 416), (0, 0, 0), True, crop=False)
        print("blobFromImage ok")
        self.net.setInput(blob)
        print("net.setInput ok")
        outs = self.net.forward(self.output_layers)
        print("net.forward ok")
        
        return outs
    
    #Creating PowerPoint
    def createPowerPoint(self) :
        self.presentation = Presentation()
        title_slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        shapes = slide.shapes
        
        return shapes

    # Showing informations on the screen
    def show(self, outs, shapes):
        print("show in")
        class_ids = []
        confidences = []
        boxes = []
        print(len(outs))
        for out in outs:
            print("for outs out in")
            for detection in out:
                scores = detection[5:]
                class_id = np.argmax(scores) 
                confidence = scores[class_id]
                if confidence > 0.6: 
                    # Object detected
                    center_x = int(detection[0] * self.width) 
                    center_y = int(detection[1] * self.height)
                    w = int(detection[2] * self.width)
                    h = int(detection[3] * self.height)
                    
                    # Rectangle coordinates
                    x = int(center_x - w / 2)
                    y = int(center_y - h / 2)
                    
                    boxes.append([x, y, w, h])
                    confidences.append(float(confidence))
                    class_ids.append(class_id)
                    
                    indexes = cv2.dnn.NMSBoxes(boxes, confidences, 0.5, 0.4) 
                    
                    font = cv2.FONT_HERSHEY_PLAIN 
        for i in range(len(boxes)):
            print("@@@@@@@@@@@@@@")
            print(len(boxes))
            if i in indexes:
                x, y, w, h = boxes[i]
                label = str(self.classes[class_ids[i]])
                print(label)
                try:
                    color = self.colors[i]
                except Exception as e:
                    print(e)
                #print(color)
                
                if label == "text":  
                    print("text in")
                    #여기에 roi를 tesseract에 보내면 됨
                    #image 인식 후 텍스트 객체? 
                    img = self.image
                    print(type(self.image))
                    print(",,,,,,,,,,,,,")
                    dst = img.copy() #원본 img copy
                    print(";;;;;;;;;;;;")
                    dst = img[y+5 : y + h + 5 , x-10 : x + w + 10]
                    roi = self.roi ######################################################
                    roi = dst
                    print("_____________________")
                    
                    #image grayscale
                    roi = cv2.cvtColor(np.asarray(roi), cv2.COLOR_BGR2GRAY)
                    roi_pil = Image.fromarray(roi)
                    roi_np = np.asarray(roi_pil)
                    s=str(i)
                    print(s)
            
                    #이미지 이진화 임계처리
                    blur= cv2.adaptiveThreshold(roi, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 61, 13)
                    print(")))))))))))))))))))))")
                    print(type(blur))
                    
                    array = blur - 18
                    image = Image.fromarray(array)
                    print("...............")
                    textImage = image.filter(ImageFilter.SHARPEN) #선명함 필터 적용
                    print("filter ok")
                    textImage.save('textImage.jpg', dpi = (400,400)) # dpi 설정
                    print("dpi save ok")
                    
                    text = pytesseract.image_to_string(textImage, lang='test',config='—psm 7 -c preserve_interword_spaces=1')
                    print(text)            
                    text = text.replace(" ", "") # 공백 제거
                    #text = " ".join(text.split())
                    print(text)
                    if text != "" :
                        if text[0].isupper() :
                            #print(text)
                            text = text.lower()
                            text = text.capitalize() # 문장 맨 첫글자만 대문자로
                            print(text)
                        else:
                            text = text.lower()
                            print(text)
                    textBox = shapes.add_textbox(Cm(x)/50, Cm(y)/50, Cm(w)/40, Cm(h)/40)
                    text_frame = textBox.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = text
                    font = run.font
                    real = w*h/650
                    if real <= 15:
                        real= 20
                    font.size = Pt(real)
                    
                    print("text out")
                    #Tesseract Text 인식 후 보정 코드 
                     
                if label == "circle":
                    print("circle in")
                    shape = shapes.add_shape(MSO_SHAPE.OVAL, Cm(x)/50, Cm(y)/50 , Cm(w)/50, Cm(h)/50)
                    shape.fill.background()
                    shape.shadow.inherit = False
                    line = shape.line
                    line.color.rgb = RGBColor(0,0,0)
                    #line.color.brightness = 0.5
                    line.width = Mm(0.5)
                    print("circle out")
                    
                if label == "rectangle":
                    print("rectangle in")
                    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x)/50, Cm(y)/50 , Cm(w)/50, Cm(h)/50)
                    shape.fill.background()
                    shape.shadow.inherit = False
                    line = shape.line
                    line.color.rgb = RGBColor(0,0,0)
                    #line.color.brightness = 0.5
                    line.width = Mm(0.5)
                    print("rectangle out")
                    
                    
                if label == "triangle":
                    print("triangle in")
                    shape = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Cm(x)/50, Cm(y)/50 , Cm(w)/50, Cm(h)/50)
                    shape.fill.background()
                    shape.shadow.inherit = False
                    line = shape.line
                    line.color.rgb = RGBColor(0,0,0)
                    #line.color.brightness = 0.5
                    line.width = Mm(0.5)
                    print("triangle out")
                    
                    
                if label == "pentagon":
                    shape = shapes.add_shape(MSO_SHAPE.REGULAR_PENTAGON, Cm(x)/85, Cm(y)/140 , Cm(w)/180, Cm(h)/180)
                    shape.fill.background()
                    shape.shadow.inherit = False
                    line = shape.line
                    line.color.rgb = RGBColor(0,0,0)
                    line.width = Mm(0.5)
                    
                if label == "arrow1":
                    #cv2.rectangle(img, (x, y), (x + w, y + h), color, 1)
                    
                    img = self.image #원본 img copy
                    gray = img[y - 9 : y + h + 9 , x-13 : x + w + 13]
                    mask =  np.zeros_like(gray)
                    mask1 = np.zeros_like(gray)
                    ret1, mask = cv2.threshold(mask, 127, 255, cv2.THRESH_BINARY_INV)
                    ret2, mask1 = cv2.threshold(mask1, 127, 255, cv2.THRESH_BINARY_INV)
                    gray = cv2.cvtColor(gray, cv2.COLOR_BGR2GRAY)
                    ret, gray = cv2.threshold(gray, 250, 255, cv2.THRESH_OTSU)
                    image_binary, contours, hierarchy = cv2.findContours(gray, cv2.RETR_TREE, cv2.CHAIN_APPROX_NONE)           
                    #drawContour on mask
                    cnt= sorted(contours, key = cv2.contourArea, reverse=True)[1]
                    cv2.drawContours(mask, cnt, -1, (0,0,0), 3)
                    #Ramer-Douglas-Peucker algorithm
                    epsilon = 0.02 * cv2.arcLength(cnt,True)
                    approx_corners = cv2.approxPolyDP(cnt, epsilon, True)
                    cv2.drawContours(mask1, approx_corners, -1, (0,0,0), 6)
                    approx_corners = sorted(np.concatenate(approx_corners).tolist())
                    approx_corners = [approx_corners[i] for i in [0, 1, 2, 3, 4, 5, 6]]
                    #무게중심
                    value_x = 0
                    value_y = 0
                    for i in approx_corners:
                        a,b = np.ravel(i)
                        value_x = value_x + a
                        value_y = value_y + b
                    
                    centerx = value_x/7
                    centery = value_y/7
                    print(approx_corners)
                    print("centerx : ", centerx)
                    print("centery : ", centery)
                    #무게중심에서 가장 가까운 두 점 계산
                    min_distancex1 = 0
                    min_distancey1 = 0
                    min_distance1 = 1000000
                    for i in approx_corners:
                        a,b = np.ravel(i)
                        distance = math.sqrt(math.pow((centerx-a),2) + math.pow((centery-b),2))
                        if distance < min_distance1:
                            min_distance1 = distance
                            min_distancex1 = a
                            min_distancey1 = b
                   
                    distance1 = np.array([min_distancex1, min_distancey1])
                    
                    print("min_distancex1 : ", min_distancex1)
                    print("min_distancey1 : ", min_distancey1)
                       
                    min_distancex2 = 0
                    min_distancey2 = 0
                    min_distance2 = 100000
                    for i in approx_corners:
                        a,b = np.ravel(i)
                        distance = math.sqrt(math.pow((centerx-a),2) + math.pow((centery-b),2))
                        if distance < min_distance2:
                            if a==min_distancex1 and b==min_distancey1:
                                continue
                            min_distance2 = distance
                            min_distancex2 = a
                            min_distancey2 = b
                            
                    distance2 = np.array([min_distancex2, min_distancey2])
                    
                    print("min_distancex2 : ", min_distancex2)
                    print("min_distancey2 : ", min_distancey2)
                    
                    center = np.array([centerx,centery])
                    
                    z1 = [i-j for i,j in zip(distance1, center)]
                    z2 = [i-j for i,j in zip(distance2, center)]
                    
                    print("z1 : ", z1)
                    print("z2 : ", z2)
                    
                    #벡터 연산
                    vector = [i+j for i,j in zip(z1,z2)]
                    a, b = np.ravel(vector)
                    
                    print("vector : ", vector)
                    
                    reference = np.array([1,0])
                    vector = np.array([a,-b])
                
                    
                    vector_size = math.sqrt(math.pow((a),2) + math.pow((b),2))
                    cosinseta = np.dot(vector, reference) / vector_size
                    print("cosinseta : ", cosinseta)
                    
        
                    degree = np.arccos(cosinseta)
                    degree = np.rad2deg(degree)
                    print("degree : ", degree)
                    
                    shape = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(x)/50, Cm(y)/50 , Cm(w)/50, Cm(h)/50)
                    #rotate clockwise
                    if 338<degree and 360>degree:
                        degree = 0
                    elif 0<degree and 22.5>degree:
                        degree = 0
                    elif 22.5<degree and 67.5>degree:
                        degree = 45
                    elif 67.5<degree and 112.5>degree:
                        degree = 90
                    elif 112.5<degree and 157.5>degree:
                        degree = 135
                    elif 157.5<degree and 202.5>degree:
                        degree = 180
                    elif 202.5<degree and 247.5>degree:
                        degree = 225
                    elif 247.5<degree and 292.5>degree:
                        degree = 270
                    elif 292.5<degree and 337.5>degree:
                        degree = 315
                    
                        
                    shape.rotation = degree
                    
                    shape.fill.background()
                    shape.shadow.inherit = False
                    line = shape.line
                    line.color.rgb = RGBColor(0,0,0)
                    line.width=Mm(0.5)
                    
                    s=str(i)
                    
                    cv2.imwrite('C:/Users/Hansung/Result/1' + s + '.jpg',mask)
                    cv2.imwrite('C:/Users/Hansung/Result/2' + s + '.jpg',mask1)
            
    
        
                #arrow 나중에 추가
                print("%%%%%%%%%%%%")
                #cv2.rectangle(img, (x, y), (x + w, y + h), (0,255,0), 2)
                print("^^^^^^^^^^^^^^^^^")
                    
        
        print("for end")
        self.presentation.save('C:/Users/hansung/test.pptx') 
        print("ppt save")
        
        
        #실제로는 욜로 인식 결과 창 안띄워도 됨, 나중에 삭제
        '''print("before show")
        cv2.imshow("Image", img)
        print("after show")
        cv2.waitKey(0)
        print('&&&&&')
        cv2.destroyAllWindows()
        print("after destroy")
        outs = self.net.forward(self.output_layers) # 이 부분 정확히 뭔지'''
        print("ooooook")
        
        #############이제 피피티 다시 전송 코드############
        return self.presentation
        
    
        
        










    
    
    
    
    
