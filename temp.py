import cv2
import numpy as np
import scipy as sp
import math
from math import *
from PIL import Image, ImageFilter
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Mm
from pptx.util import Cm
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import matplotlib.pyplot as plt
import matplotlib as mpl

#def sobel_filters(img):
#    Kx = np.array([[-1,0,1],[-2,0,2],[-1,0,1]], np.float32)
#    Ky = np.array([[1,2,1],[0,0,0],[-1,-2,-1]], np.float32)
#    
#    Ix = ndimage.filters.convolve(img, Kx)
#    Iy = ndimage.filters.convolve(img, Ky)
#    
#    G = np.hypot(Ix, Iy)
#    G = G / G.max() * 255
#    theta = np.arctan2(Iy, Ix)
#    
#    return (theta*180/3.14)

def get_gradient(image, kernel_size):
    
    grad_x = cv2.Sobel(image, cv2.CV_32F, 1, 0 ,ksize=kernel_size)
    grad_y = cv2.Sobel(image, cv2.CV_32F, 0, 1 ,ksize=kernel_size)
    
    grad = grad_x + 1j*grad_y
    
    return grad

# Load Yolo
net = cv2.dnn.readNet("C:/Users/Hansung/Desktop/project/Yolo/darkflow/bin/yolo-obj_40000.weights", "C:/Users/Hansung/Desktop/project/Yolo/darkflow/cfg/yolo-obj.cfg")
classes = []                ##weight file
mpl.rcParams["figure.dpi"] = 300
with open("C:/Users/Hansung/Yolo_mark-master/x64/Release/data/obj.names", "r") as f:
    classes = [line.strip() for line in f.readlines()]
layer_names = net.getLayerNames()
output_layers = [layer_names[i[0] - 1] for i in net.getUnconnectedOutLayers()]
colors = np.random.uniform(0, 255, size=(len(classes), 3))
roi=0
# Loading image
img = cv2.imread("C:/Users/Hansung/Desktop/project/all2.jpg")
#img = cv2.resize(img, None, fx=0.6, fy=0.6, interpolation=cv2.INTER_LINEAR)
img = cv2.resize(img,(1280,960))
height, width, channels = img.shape

# Detecting objects
blob = cv2.dnn.blobFromImage(img, 0.00392, (416, 416), (0, 0, 0), True, crop=False)
net.setInput(blob)
outs = net.forward(output_layers)

#Creating PowerPoint
presentation = Presentation()
title_slide_layout = presentation.slide_layouts[6]
slide = presentation.slides.add_slide(title_slide_layout)
shapes = slide.shapes

# Showing informations on the screen
class_ids = []
confidences = []
boxes = []
list = []
for out in outs:
    for detection in out:
        scores = detection[5:]
        class_id = np.argmax(scores)
        confidence = scores[class_id]
        if confidence > 0.4:
            # Object detected
            center_x = int(detection[0] * width)
            center_y = int(detection[1] * height)
            w = int(detection[2] * width)
            h = int(detection[3] * height)
            
            # Rectangle coordinates
            x = int(center_x - w / 2)
            y = int(center_y - h / 2)
            
            boxes.append([x, y, w, h])

            confidences.append(float(confidence))
            class_ids.append(class_id)
            
indexes = cv2.dnn.NMSBoxes(boxes, confidences, 0.5, 0.4)
font = cv2.FONT_HERSHEY_PLAIN
            
for i in range(len(boxes)):
    
    if i in indexes:
        x, y, w, h = boxes[i]
        label = str(classes[class_ids[i]])
        color = colors[1]
        cv2.rectangle(img, (x, y), (x + w, y + h), color, 2)
        cv2.putText(img, label, (x, y + 30), font, 3, color, 3)
        if label == "text":  
            #여기에 roi를 tesseract에 보내면 됨
            dst = img.copy() #원본 img copy
            dst = img[y - 7 : y + h + 7 , x-10 : x + w + 10]
            roi = dst
            #image grayscale
            roi = cv2.cvtColor(np.asarray(roi), cv2.COLOR_BGR2GRAY)
            roi_pil = Image.fromarray(roi)
            roi_np = np.asarray(roi_pil)
            s=str(i)
            
                #Otsu threshold & Gaussian blur
            #ret1, th1 = cv2.threshold(roi, 0, 255, cv2.THRESH_BINARY+cv2.THRESH_OTSU)
            #blur = cv2.GaussianBlur(roi, (1,1), 0)
            
                #AdaptiveThreshold & Gaussian blur
            blur= cv2.adaptiveThreshold(roi, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 61, 13)
            
                #erode & dilate
            #kernel = np.ones((3,3), np.uint8)
            #blur= cv2.erode(blur, kernel, 1)
            #kernel = np.ones((3,3), np.uint8)
            #blur = cv2.dilate(blur, kernel, 1)
            
            #result = cv2.resize(result, None, fx=0.2, fy=0.2)
            cv2.imwrite('C:/Users/Hansung/Desktop/project/7' + s + '.jpg',blur) #label == text일때 이미지저장
            #dpi 조절
            source_image = 'C:/Users/Hansung/Desktop/project/7' + s + '.jpg'
            image = Image.open(source_image)
            image = image.filter(ImageFilter.SHARPEN)
            image.save('C:/Users/Hansung/Desktop/project/7' + s + '.jpg', dpi = (400,400))
                     
        if label == "circle":
            
            shape = shapes.add_shape(MSO_SHAPE.OVAL, Cm(x)/50, Cm(y)/50 , Cm(w)/50, Cm(h)/50)
            shape.fill.background()
            shape.shadow.inherit = False
            line = shape.line
            line.color.rgb = RGBColor(0,0,0)
            #line.color.brightness = 0.5
            #line.width = Mm(1)
        
            
        if label == "rectangle":
            shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x)/50, Cm(y)/50 , Cm(w)/50, Cm(h)/50)
            shape.fill.background()
            shape.shadow.inherit = False
            line = shape.line
            line.color.rgb = RGBColor(0,0,0)
            #line.color.brightness = 0.5
           # line.width = Mm(1)
            
            
        if label == "triangle":
            shape = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Cm(x)/50, Cm(y)/50 , Cm(w)/50, Cm(h)/50)
            shape.fill.background()
            shape.shadow.inherit = False
            line = shape.line
            line.color.rgb = RGBColor(0,0,0)
            #line.color.brightness = 0.5
           # line.width = Mm(1)
            
            
        if label == "pentagon":
            shape = shapes.add_shape(MSO_SHAPE.REGULAR_PENTAGON, Cm(x)/85, Cm(y)/140 , Cm(w)/180, Cm(h)/180)
            shape.fill.background()
            shape.shadow.inherit = False
            line = shape.line
            line.color.rgb = RGBColor(0,0,0)
            #line.color.brightness = 0.5
           # line.width = Mm(1)
        
        #수정해야되는부분
        if label == "arrow1":
            '''
            img = img.copy() #원본 img copy
            gray = img[y - 7 : y + h + 7 , x-10 : x + w + 10]
            
            gray_image = cv2.cvtColor(gray, cv2.COLOR_BGR2GRAY)
            gray_image = cv2.GaussianBlur(gray_image, (3,3), 3)
            
            
            #Convex Hull
            ret, img_binary = cv2.threshold(gray_image, 127, 255, 0)
            img_binary, contours, hierarchy = cv2.findContours(img_binary, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
            
            
            for cnt in contours:
                hull=cv2.convexHull(cnt)
                cv2.drawContours(gray,[hull], 0, (255,255,0), 6)

            for cnt in contours:
                hull = cv2.convexHull(cnt, returnPoints = False)
                defects = cv2.convexityDefects(cnt, hull)

                for i in range(defects.shape[0]):
                    s,e,f,d = defects[i,0]
                    start = tuple(cnt[s][0])
                    end = tuple(cnt[e][0])
                    far = tuple(cnt[f][0])
                    
                    print(d)
                    
                    cv2.circle(gray, far, 4, (0,255,0),-1)                
            s=str(i)
            cv2.imwrite('C:/Users/Hansung/Desktop/project/1' + s + '.jpg',gray)
            '''
            
        if label == "arrow2":
            #img = img.copy() #원본 img copy
            #gray = img[y - 7 : y + h + 7 , x-10 : x + w + 10]
            
            #gray_image = cv2.cvtColor(gray, cv2.COLOR_BGR2GRAY)
            #gray_image = cv2.GaussianBlur(gray_image, (11,11), 3)
            
            
            #Convex Hull
            #ret, img_binary = cv2.threshold(gray_image, 127, 255, 0)
            #img_binary, contours, hierarchy = cv2.findContours(img_binary, cv2.RETR_CCOMP, cv2.CHAIN_APPROX_SIMPLE)
            
 
            
            #for cnt in contours:
             #   hull = cv2.convexHull(cnt)
              #  cv2.drawContours(gray,[hull], 0, (255,0,255), 6)
                

            for cnt in contours:
                hull = cv2.convexHull(cnt, returnPoints = False)
                defects = cv2.convexityDefects(cnt, hull)

                for i in range(defects.shape[0]):
                    s,e,f,d = defects[i,0]
                    start = tuple(cnt[s][0])
                    end = tuple(cnt[e][0])
                    far = tuple(cnt[f][0])
                    
                    print(d)
                    
                    cv2.circle(gray, far, 5, (0,255,0),-1)
           
            s=str(i)
            cv2.imwrite('C:/Users/Hansung/Desktop/project/1' + s + '.jpg',gray)
            
            


    
presentation.save('C:/Users/Hansung/Desktop/project/test.pptx')

cv2.imshow("Image", img)
cv2.waitKey(0)
cv2.destroyAllWindows()
outs = net.forward(output_layers)